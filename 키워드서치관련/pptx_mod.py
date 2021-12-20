import zipfile
from collections import OrderedDict
from pptx import Presentation
try:
    from xml.etree.cElementTree import XML
except ImportError:
    from xml.etree.ElementTree import XML
from timeit import default_timer as timer
from datetime import timedelta

def xml(full_path):
    document = zipfile.ZipFile(full_path)
    nums = []
    value_list = []
    for d in document.namelist():
        if d.startswith("ppt/slides/slide"):
            nums.append(int(d[len("ppt/slides/slide"):-4]))
    s_format = "ppt/slides/slide%s.xml"
    slide_name_list = [s_format % x for x in sorted(nums)]
    for slide in slide_name_list:
        xml_content = document.read(slide)
        tree = XML(xml_content)
        if tree.tag == '{http://purl.oclc.org/ooxml/presentationml/main}sld':
            NAMESPACE = '{http://purl.oclc.org/ooxml/drawingml/main}'
        else:
            NAMESPACE = '{http://schemas.openxmlformats.org/drawingml/2006/main}'
        TEXT = NAMESPACE + 't'
        for node in tree.iter(TEXT):
            value = str(node.text).upper()
            value_list.append(value)
    value_list = list(OrderedDict.fromkeys(value_list))
    txt = ','.join(value_list)
    print(txt)
    document.close()

def pptx_module(full_path):
    prs = Presentation(full_path)
    result = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                result.append(paragraph.text)
    result = list(OrderedDict.fromkeys(result))
    txt = ','.join(result)
    print(txt)

if __name__ == '__main__':
    start = timer()

    full_path = r'X:\2021년\03월\20210329 티피아이인사이트\복원\0000\C_삭제 자료\Extra Found Files\Document\Microsoft PowerPoint 2007 XML Document\더앤컴퍼니9월 월례회.pptx'
    xml(full_path)

    end = timer()
    print(timedelta(seconds=end - start))

